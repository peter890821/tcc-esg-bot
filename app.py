import io
import os
import re
import tempfile
from dataclasses import dataclass
from typing import List, Optional, Tuple

import pandas as pd
import streamlit as st
from docx import Document
from pptx import Presentation
from pypdf import PdfReader

try:
    import google.generativeai as genai
except Exception:  # pragma: no cover
    genai = None

try:
    import chromadb
    from chromadb.utils import embedding_functions
except Exception:  # pragma: no cover
    chromadb = None

try:
    from unstructured.partition.pdf import partition_pdf
    from bs4 import BeautifulSoup
except Exception:  # pragma: no cover
    partition_pdf = None
    BeautifulSoup = None


DEFAULT_CSV_PATH = "MSCI_Methodology_Full_KB.csv"
KB_DIRECTORY = "kb"  # 專門放長期知識庫檔案的資料夾
REQUIRED_COLUMNS = {"text_content", "source_file", "doc_type"}


@dataclass
class RetrievedChunk:
    text_content: str
    source_file: str
    doc_type: str
    score: int
    year: Optional[int] = None
    version: Optional[str] = None
    page_number: Optional[int] = None


def _safe_str(x) -> str:
    return "" if x is None else str(x)


def sanitize_text(text: str) -> str:
    """
    清洗字串，移除無法被 UTF-8 編碼的 surrogate 字元。
    這通常發生在處理複雜 PDF 提取出的損壞字元。
    """
    if not isinstance(text, str):
        return str(text)
    # 使用 'ignore' 拋棄無法編碼的字元，再解碼回來
    return text.encode('utf-8', 'ignore').decode('utf-8', 'ignore')


def tokenize_question(question: str) -> List[str]:
    """
    Very simple tokenization:
    - Extract alphanumeric "words" (English/numbers)
    - Extract CJK sequences (Chinese/Japanese/Korean characters)
    """
    q = question.strip().lower()
    if not q:
        return []

    tokens: List[str] = []
    tokens += re.findall(r"[a-z0-9]+", q)
    tokens += re.findall(r"[\u4e00-\u9fff]+", q)

    # Keep unique tokens, preserve order
    seen = set()
    deduped = []
    for t in tokens:
        if t and t not in seen:
            seen.add(t)
            deduped.append(t)
    return deduped


def simple_retrieve_topk(df: pd.DataFrame, question: str, k: int = 5) -> List[RetrievedChunk]:
    """
    Simple string matching retrieval (fallback when vector search unavailable):
    score = sum(count(token in text_content)) + (bonus if full question is substring)
    """
    tokens = tokenize_question(question)
    if not tokens:
        return []

    q_lower = question.strip().lower()

    chunks: List[RetrievedChunk] = []
    for idx, row in df.iterrows():
        text = _safe_str(row.get("text_content"))
        if not text:
            continue
        text_lower = text.lower()

        score = 0
        for t in tokens:
            score += text_lower.count(t)
        if q_lower and q_lower in text_lower:
            score += 5

        if score > 0:
            chunks.append(
                RetrievedChunk(
                    text_content=text,
                    source_file=_safe_str(row.get("source_file")),
                    doc_type=_safe_str(row.get("doc_type")),
                    score=score,
                    year=row.get("year") if pd.notna(row.get("year")) else None,
                    version=_safe_str(row.get("version")) if row.get("version") else None,
                    page_number=row.get("page_number") if pd.notna(row.get("page_number")) else None,
                )
            )

    chunks.sort(key=lambda c: c.score, reverse=True)
    return chunks[:k]


@st.cache_resource
def init_vector_db(df: pd.DataFrame, api_key: str) -> Optional[object]:
    """
    初始化向量資料庫，將知識庫嵌入
    使用 ChromaDB + Gemini Embeddings
    """
    if chromadb is None:
        st.warning("ChromaDB 未安裝，將使用簡易字串檢索。建議執行：pip install chromadb")
        return None
    
    try:
        client = chromadb.Client()
        
        # 使用 Gemini Embedding API
        gemini_ef = embedding_functions.GoogleGenerativeAiEmbeddingFunction(
            api_key=api_key,
            model_name="models/text-embedding-004"
        )
        
        # 建立 collection（若已存在則刪除重建）
        try:
            client.delete_collection(name="tcc_esg_kb")
        except:
            pass
        
        collection = client.create_collection(
            name="tcc_esg_kb",
            embedding_function=gemini_ef,
            metadata={"hnsw:space": "cosine"}
        )
        
        # 批量嵌入文檔
        valid_texts = []
        valid_metadatas = []
        valid_ids = []
        
        for idx, row in df.iterrows():
            text = _safe_str(row.get("text_content"))
            if text.strip():  # 只加入非空文本
                valid_texts.append(text)
                valid_ids.append(f"chunk_{idx}")
                valid_metadatas.append({
                    "source_file": _safe_str(row.get("source_file")),
                    "doc_type": _safe_str(row.get("doc_type")),
                    "year": int(row.get("year")) if pd.notna(row.get("year")) else 0,
                    "idx": int(idx)
                })
        
        if valid_texts:
            collection.add(
                ids=valid_ids,
                documents=valid_texts,
                metadatas=valid_metadatas
            )
        
        return collection
    
    except Exception as e:
        st.error(f"向量資料庫初始化失敗：{e}。將使用簡易檢索。")
        return None


def extract_chinese_terms(text: str) -> List[str]:
    """提取中文重要術語（3字以上）"""
    terms = re.findall(r'[\u4e00-\u9fff]{3,}', text)
    return list(set(terms))


def hybrid_retrieve_topk(
    collection: Optional[object],
    df: pd.DataFrame,
    question: str,
    k: int = 5
) -> List[RetrievedChunk]:
    """
    混合檢索：向量語義搜尋 (70%) + 關鍵字精確匹配 (30%)
    適用於 ESG 領域的專業術語檢索
    """
    # Fallback to simple search if vector DB unavailable
    if collection is None:
        return simple_retrieve_topk(df, question, k)
    
    try:
        # 1. 向量搜尋（語義理解）
        results = collection.query(
            query_texts=[question],
            n_results=min(k * 2, len(df))  # 取雙倍候選
        )
        
        vector_scores = {}
        if results['ids'] and results['ids'][0]:
            for i, chunk_id in enumerate(results['ids'][0]):
                idx = int(chunk_id.split('_')[1])
                # 距離轉為相似度分數 (0-1)
                distance = results['distances'][0][i]
                similarity = max(0, 1.0 - distance)
                vector_scores[idx] = similarity
        
        # 2. 關鍵字精確匹配（專有名詞）
        keyword_patterns = [
            r'ISO\s*\d+',
            r'IFRS\s*S\d+',
            r'GRI\s*\d+',
            r'Scope\s*[123]',
            r'TCFD',
            r'SASB',
            r'CBAM',
            r'MSCI'
        ]
        
        keyword_scores = {}
        q_lower = question.lower()
        important_terms = extract_chinese_terms(question)
        
        for idx, row in df.iterrows():
            text = _safe_str(row.get("text_content"))
            if not text:
                continue
            
            text_lower = text.lower()
            keyword_score = 0.0
            
            # 檢查專有名詞 (高權重)
            for pattern in keyword_patterns:
                if re.search(pattern, question, re.IGNORECASE):
                    matches = len(re.findall(pattern, text, re.IGNORECASE))
                    keyword_score += matches * 0.3
            
            # 檢查中文重要術語 (中權重)
            for term in important_terms:
                if term in text:
                    keyword_score += text.count(term) * 0.2
            
            # 檢查完整問題匹配 (高權重)
            if len(q_lower) > 5 and q_lower in text_lower:
                keyword_score += 0.5
            
            if keyword_score > 0:
                keyword_scores[idx] = min(keyword_score, 1.0)  # 限制最大值為1
        
        # 3. 混合排序（向量 70% + 關鍵字 30%）
        combined_scores = {}
        
        # 合併所有候選
        all_indices = set(list(vector_scores.keys()) + list(keyword_scores.keys()))
        
        for idx in all_indices:
            vector_score = vector_scores.get(idx, 0.0)
            keyword_score = keyword_scores.get(idx, 0.0)
            combined_scores[idx] = vector_score * 0.7 + keyword_score * 0.3
        
        # 排序並取 Top-K
        sorted_indices = sorted(
            combined_scores.items(),
            key=lambda x: x[1],
            reverse=True
        )[:k]
        
        # 構建結果
        results = []
        for idx, score in sorted_indices:
            row = df.iloc[idx]
            results.append(
                RetrievedChunk(
                    text_content=_safe_str(row.get("text_content")),
                    source_file=_safe_str(row.get("source_file")),
                    doc_type=_safe_str(row.get("doc_type")),
                    score=int(score * 100),  # 轉為百分比
                    year=int(row.get("year")) if pd.notna(row.get("year")) else None,
                    version=_safe_str(row.get("version")) if row.get("version") else None,
                    page_number=int(row.get("page_number")) if pd.notna(row.get("page_number")) else None,
                )
            )
        
        return results
    
    except Exception as e:
        st.warning(f"向量檢索失敗：{e}。切換至簡易檢索。")
        return simple_retrieve_topk(df, question, k)


def html_table_to_markdown(html: str) -> str:
    """將 HTML 表格轉為 Markdown（保留結構，AI 更易理解）"""
    if BeautifulSoup is None:
        return html
    
    try:
        soup = BeautifulSoup(html, 'html.parser')
        table = soup.find('table')
        
        if not table:
            return html
        
        rows = table.find_all('tr')
        if not rows:
            return html
        
        md_lines = []
        
        for i, row in enumerate(rows):
            cells = row.find_all(['th', 'td'])
            if not cells:
                continue
            
            # 提取單元格文字
            cell_texts = [c.get_text(strip=True) for c in cells]
            md_lines.append('| ' + ' | '.join(cell_texts) + ' |')
            
            # 第一行後添加表頭分隔線
            if i == 0:
                md_lines.append('| ' + ' | '.join(['---'] * len(cells)) + ' |')
        
        return '\n'.join(md_lines)
    
    except Exception:
        return html


def _extract_text_from_pdf(file_bytes: bytes, filename: str = "document.pdf") -> str:
    """
    使用 Unstructured 解析 PDF，保留表格結構
    若 Unstructured 不可用，回退至基本解析
    """
    # 嘗試使用 Unstructured（表格感知）
    if partition_pdf is not None:
        try:
            # 暫存檔案（Unstructured 需要檔案路徑）
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp:
                tmp.write(file_bytes)
                tmp_path = tmp.name
            
            try:
                # 分區解析（自動檢測表格）
                elements = partition_pdf(
                    filename=tmp_path,
                    strategy="hi_res",  # 高解析度，啟用表格檢測
                    infer_table_structure=True,
                    extract_images_in_pdf=False
                )
                
                # 分類處理
                text_parts = []
                for elem in elements:
                    if elem.category == "Table":
                        # 表格轉 Markdown 格式
                        table_html = elem.metadata.text_as_html if hasattr(elem.metadata, 'text_as_html') else str(elem)
                        table_md = html_table_to_markdown(table_html)
                        text_parts.append(f"\n[表格]\n{table_md}\n")
                    else:
                        text_parts.append(elem.text)
                
                return sanitize_text("\n".join(text_parts))
            
            finally:
                # 清理暫存檔
                try:
                    os.unlink(tmp_path)
                except:
                    pass
        
        except Exception as e:
            # Unstructured 失敗，回退至基本解析
            st.warning(f"高級 PDF 解析失敗（{filename}）：{e}。使用基本解析。")
    
    # 基本解析（Fallback）
    reader = PdfReader(io.BytesIO(file_bytes))
    texts: List[str] = []
    for page in reader.pages:
        try:
            page_text = page.extract_text() or ""
        except Exception:
            page_text = ""
        if page_text:
            texts.append(page_text)
    return sanitize_text("\n\n".join(texts))


def auto_extract_year(filename: str, content: str = "") -> Optional[int]:
    """
    從檔名或內容自動提取年份
    
    支援格式：
    - 西元年：2024, 2023
    - 民國年：民國 113 年 → 2024
    - ISO 日期：2024-01-15
    """
    # 1. 檔名優先（例如：TCC_ESG_Report_2024.pdf）
    year_match = re.search(r'20\d{2}', filename)
    if year_match:
        year = int(year_match.group())
        # 驗證合理性（2000-2050）
        if 2000 <= year <= 2050:
            return year
    
    # 2. 內容中查找（只檢查前 2000 字元）
    if content:
        sample = content[:2000]
        
        # 民國年（例如：「民國 113 年」 → 2024）
        roc_match = re.search(r'民國\s*(\d{3})', sample)
        if roc_match:
            roc_year = int(roc_match.group(1))
            if 100 <= roc_year <= 150:  # 合理範圍
                return roc_year + 1911
        
        # ISO 日期格式（例如：2024-01-15）
        iso_match = re.search(r'20\d{2}-\d{2}-\d{2}', sample)
        if iso_match:
            year = int(iso_match.group()[:4])
            if 2000 <= year <= 2050:
                return year
        
        # 純西元年（例如：「2024 年度報告」）
        year_match = re.search(r'(20\d{2})\s*年', sample)
        if year_match:
            year = int(year_match.group(1))
            if 2000 <= year <= 2050:
                return year
    
    return None


def extract_version(filename: str) -> Optional[str]:
    """
    從檔名提取版本號
    
    支援格式：
    - v1.0, V2.3
    - version_1.0
    - _final, _draft
    """
    # 版本號格式（v1.0, V2.3）
    version_match = re.search(r'[vV](\d+\.\d+)', filename)
    if version_match:
        return f"v{version_match.group(1)}"
    
    # version_xxx
    version_match = re.search(r'version[_\s]*(\d+(?:\.\d+)?)', filename, re.IGNORECASE)
    if version_match:
        return f"v{version_match.group(1)}"
    
    # 狀態標記
    if '_final' in filename.lower() or '-final' in filename.lower():
        return "final"
    if '_draft' in filename.lower() or '-draft' in filename.lower():
        return "draft"
    
    return None


def _extract_text_from_docx(file_bytes: bytes) -> str:
    doc = Document(io.BytesIO(file_bytes))
    paras = [p.text for p in doc.paragraphs if p.text]
    return sanitize_text("\n".join(paras))


def _extract_text_from_pptx(file_bytes: bytes) -> str:
    pres = Presentation(io.BytesIO(file_bytes))
    texts: List[str] = []
    for slide in pres.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
    return sanitize_text("\n\n".join(texts))


@st.cache_data(show_spinner=False)
def load_kb_from_bytes(file_bytes: bytes, filename: str) -> pd.DataFrame:
    """
    Load uploaded knowledge base file into a standardized DataFrame
    with columns: text_content, source_file, doc_type, year, version.

    支援類型：CSV、Excel、PDF、Word、PPT。
    - 對於非 CSV/Excel，我們會把整個檔案文字當成一個段落。
    - 自動提取年份與版本 metadata
    """
    name = filename or "uploaded_file"
    lower = name.lower()
    ext = ""
    if "." in lower:
        ext = lower.rsplit(".", 1)[-1]

    buffer = io.BytesIO(file_bytes)

    # Structured: CSV / Excel
    if ext in {"csv", "xlsx"}:
        if ext == "xlsx":
            table = pd.read_excel(buffer)
        else:
            table = pd.read_csv(buffer)

        cols = set(str(c) for c in table.columns)
        # 若已經有標準欄位，就直接使用
        if REQUIRED_COLUMNS.issubset(cols):
            df = table.copy()
            for col in REQUIRED_COLUMNS:
                df[col] = df[col].fillna("").astype(str)
            
            # 補充 metadata 欄位（若不存在）
            if "year" not in df.columns:
                df["year"] = auto_extract_year(name, "")
            if "version" not in df.columns:
                df["version"] = extract_version(name)
            
            return df[["text_content", "source_file", "doc_type", "year", "version"]]

        # 否則：把整張表轉成一段文字
        text_repr = table.astype(str).to_csv(index=False)
        return pd.DataFrame(
            [
                {
                    "text_content": text_repr,
                    "source_file": name,
                    "doc_type": ext or "table",
                    "year": auto_extract_year(name, text_repr),
                    "version": extract_version(name),
                }
            ]
        )

    # Unstructured: PDF / DOCX / PPTX
    if ext == "pdf":
        text = _extract_text_from_pdf(file_bytes, filename=name)
    elif ext in {"docx", "doc"}:
        text = _extract_text_from_docx(file_bytes)
    elif ext in {"pptx", "ppt"}:
        text = _extract_text_from_pptx(file_bytes)
    else:
        # Fallback：當作一般文字檔
        try:
            text = buffer.read().decode("utf-8")
        except Exception:
            text = ""

    return pd.DataFrame(
        [
            {
                "text_content": text,
                "source_file": name,
                "doc_type": ext or "file",
                "year": auto_extract_year(name, text),
                "version": extract_version(name),
            }
        ]
    )


@st.cache_data(show_spinner=False)
def load_csv_from_path(path: str) -> pd.DataFrame:
    return pd.read_csv(path)


def get_kb_dataframe(uploaded_files) -> Tuple[Optional[pd.DataFrame], Optional[str]]:
    """
    Returns (df, error_message). error_message is None when success.

    知識來源優先順序：
    1. 專案內的 kb/ 資料夾（長期知識庫，啟動時自動載入）
    2. 使用者在網頁上臨時上傳的檔案
    3. 專案根目錄下的預設 CSV：MSCI_Methodology_Full_KB.csv
    """
    dfs: List[pd.DataFrame] = []

    # 1) 掃描本地 kb/ 資料夾（遞迴掃描所有子資料夾）
    if os.path.isdir(KB_DIRECTORY):
        for root, _, files in os.walk(KB_DIRECTORY):
            for name in files:
                lower = name.lower()
                if not lower.endswith((".csv", ".xlsx", ".pdf", ".docx", ".pptx")):
                    continue
                path = os.path.join(root, name)
                # 保留相對於 kb/ 的完整路徑（包含子資料夾），例如：環境相關/report.pdf
                rel_path = os.path.relpath(path, KB_DIRECTORY)
                try:
                    with open(path, "rb") as f:
                        file_bytes = f.read()
                    df_part = load_kb_from_bytes(file_bytes, rel_path)
                    dfs.append(df_part)
                except Exception:
                    # 若單一檔案失敗，不影響整體，可加日後 logging
                    continue

    # 2) 網頁上臨時上傳的檔案
    if uploaded_files:
        try:
            for f in uploaded_files:
                file_bytes = f.getvalue()
                df_part = load_kb_from_bytes(file_bytes, f.name)
                dfs.append(df_part)
        except Exception as e:
            return None, f"讀取上傳檔案失敗：{e}"

    # 3) 若前兩者都沒有資料，試著載入預設 CSV
    if not dfs:
        try:
            df = load_csv_from_path(DEFAULT_CSV_PATH)
        except FileNotFoundError:
            return None, (
                f"找不到任何知識庫資料。\n"
                f"- 若要使用預設檔案，請將 {DEFAULT_CSV_PATH} 放在專案根目錄。\n"
                f"- 或建立 `{KB_DIRECTORY}` 資料夾，放入 CSV/Excel/PDF/Word/PPT 檔案。\n"
                f"- 或直接在左側上傳 ESG 知識庫檔案。"
            )
        except pd.errors.EmptyDataError:
            return None, "預設 CSV 檔案是空的或格式不正確。"
        except UnicodeDecodeError:
            return None, "預設 CSV 編碼讀取失敗。請嘗試另存為 UTF-8 後重新放置。"
        except Exception as e:
            return None, f"讀取預設 CSV 失敗：{e}"

        missing = REQUIRED_COLUMNS - set(df.columns)
        if missing:
            return None, f"預設 CSV 缺少必要欄位：{', '.join(sorted(missing))}。需要欄位：{', '.join(sorted(REQUIRED_COLUMNS))}"

        for col in REQUIRED_COLUMNS:
            df[col] = df[col].fillna("").astype(str)

        return df[list(REQUIRED_COLUMNS)], None

    # 合併來自 kb/ 與上傳的所有資料
    df_all = pd.concat(dfs, ignore_index=True)
    for col in REQUIRED_COLUMNS:
        if col not in df_all.columns:
            df_all[col] = ""
        df_all[col] = df_all[col].fillna("").astype(str)

    return df_all[list(REQUIRED_COLUMNS)], None


def build_prompt(context: str, question: str) -> str:
    """
    構建專業 ESG 顧問 Prompt，整合框架對照與結構化輸出
    """
    # 清洗輸入，確保沒有 surrogate 字元
    clean_context = sanitize_text(context)
    clean_question = sanitize_text(question)
    return (
        "你現在是台泥集團 (TCC) 的首席永續策略顧問與 ESG 分析專家。\n\n"
        
        "【核心任務】\n"
        f"根據以下背景資料回答使用者問題：\n{clean_context}\n\n"
        f"使用者問題：{clean_question}\n\n"
        
        "【回答要求】\n"
        "1. 使用繁體中文，專業、結構化的方式回答\n"
        "2. 明確引用相關 ESG 框架與條文編號（若適用）\n"
        "3. 提供具體、可執行的建議\n\n"
        
        "【重點框架對照】\n"
        "在回答時，請主動識別並引用以下相關框架標準：\n"
        "- **IFRS S1/S2**：永續揭露準則（一般要求與氣候相關揭露）\n"
        "- **GRI Standards 2021**：全球報告倡議組織標準\n"
        "- **SASB (EM-CM)**：永續會計準則（建材產業）\n"
        "- **TCFD**：氣候相關財務揭露\n"
        "- **EU CBAM**：歐盟碳邊境調整機制\n"
        "- **MSCI ESG**：MSCI ESG 評等方法論\n\n"
        
        "【回答格式】\n"
        "請按照以下結構組織你的回答：\n"
        "1️⃣ **直接回答**\n"
        "   - 簡潔回應核心問題\n\n"
        "2️⃣ **框架引用**（若適用）\n"
        "   - 例如：根據 IFRS S2 第 14 條...\n"
        "   - 例如：符合 GRI 305-1（直接溫室氣體排放）...\n"
        "   - 例如：依據 SASB EM-CM-110a.1（能源管理）...\n\n"
        "3️⃣ **TCC 具體建議**\n"
        "   - 針對台泥集團的實務建議\n\n"
        "4️⃣ **參考依據**\n"
        "   - 列出關鍵參考來源檔案或段落\n"
    )


def generate_with_gemini(api_key: str, prompt: str) -> str:
    if genai is None:
        raise RuntimeError("找不到 google-generativeai 套件。請確認已安裝 requirements.txt 內的依賴。")

    genai.configure(api_key=api_key)

    # 先查出目前這個 API Key 可用的模型
    try:
        available_models = list(genai.list_models())
    except Exception as e:  # pragma: no cover
        raise RuntimeError(
            "Gemini 呼叫失敗：無法取得可用模型清單，請確認：\n"
            "1. 這支 API Key 是在 Google AI Studio 產生的，而不是 Google Cloud / 其他服務。\n"
            "2. 已在對應的帳號啟用 Gemini API，且網路可以連到 Google。\n"
            f"詳細錯誤：{e}"
        ) from e

    text_models = [
        m
        for m in available_models
        if hasattr(m, "supported_generation_methods")
        and "generateContent" in getattr(m, "supported_generation_methods", [])
    ]

    if not text_models:
        raise RuntimeError(
            "Gemini 呼叫失敗：這支 API Key 名下目前沒有支援 generateContent 的模型可用。\n"
            "請到 Google AI Studio 的 Models / API 頁面，確認帳號已開通 Gemini 1.5（例如 gemini-1.5-flash）後再試一次。"
        )

    # 優先挑選名稱中包含 1.5 的模型，其次任意可用模型
    preferred = [m for m in text_models if "1.5" in getattr(m, "name", "")]
    chosen = (preferred or text_models)[0]
    model_name = getattr(chosen, "name", "gemini-1.5-flash")

    try:
        model = genai.GenerativeModel(model_name)
        # 在發送前再次確保 prompt 是乾淨的 UTF-8
        clean_prompt = sanitize_text(prompt)
        resp = model.generate_content(clean_prompt)
    except Exception as e:  # pragma: no cover
        raise RuntimeError(
            "Gemini 呼叫失敗：雖然成功找到可用模型 "
            f"`{model_name}`，但在呼叫 generateContent 時出現錯誤：{e}\n"
            "請到 Google AI Studio 測試同一支 API Key 是否可以正常呼叫相同模型。"
        ) from e

    text = getattr(resp, "text", None)
    if not text:
        # Fallback for some response shapes
        try:
            text = resp.candidates[0].content.parts[0].text  # type: ignore[attr-defined]
        except Exception:
            text = ""

    if not text.strip():
        raise RuntimeError("Gemini 未回傳有效文字內容，請稍後再試或確認 API Key/配額。")

    return text.strip()


def main() -> None:
    st.set_page_config(page_title="TCC ESG Intelligent Knowledge Base", layout="wide")

    st.title("🌿 台泥 (TCC) 企業 ESG 智能知識庫")
    st.caption("Powered by Google Gemini 1.5 Pro & MSCI Methodology")

    with st.sidebar:
        st.header("⚙️ 設定 (Settings)")
        api_key = st.text_input("Google Gemini API Key", type="password", value="")
        uploaded_files = st.file_uploader(
            "上傳 ESG 知識庫檔案（可多選：CSV / Excel / PDF / Word / PPT）",
            type=["csv", "xlsx", "pdf", "docx", "pptx"],
            accept_multiple_files=True,
        )

    df, kb_err = get_kb_dataframe(uploaded_files)
    
    # Metadata filtering UI (after loading KB)
    with st.sidebar:
        if df is not None and not kb_err:
            st.divider()
            st.header("📅 資料篩選 (Filters)")
            
            # Year filter
            if 'year' in df.columns:
                available_years = sorted([int(y) for y in df['year'].dropna().unique()])
                if available_years:
                    selected_years = st.multiselect(
                        "🗓️ 文件年份",
                        options=available_years,
                        default=[],
                        help="選擇特定年份的文件（留空=全部）"
                    )
                    
                    # Apply year filter
                    if selected_years:
                        df = df[df['year'].isin(selected_years)]
                        st.success(f"✓ 已篩選 {len(selected_years)} 個年份")
            
            # Document type filter
            if 'doc_type' in df.columns:
                available_types = sorted(df['doc_type'].dropna().unique())
                if available_types:
                    selected_doc_types = st.multiselect(
                        "📂 文件類型",
                        options=available_types,
                        default=[],
                        help="選擇特定類型的文件（留空=全部）"
                    )
                    
                    # Apply doc type filter
                    if selected_doc_types:
                        df = df[df['doc_type'].isin(selected_doc_types)]
                        st.success(f"✓ 已篩選 {len(selected_doc_types)} 種類型")
    if kb_err:
        st.error(kb_err)

    if not api_key.strip():
        st.info("請先在側邊欄輸入 **Google Gemini API Key** 才能開始分析。")

    question = st.text_area(
        "請輸入你的問題",
        placeholder="例如：MSCI 對於環境管理（E）在評等方法論中通常會如何衡量？",
        disabled=not api_key.strip(),
    )

    run = st.button("開始分析", type="primary", disabled=(not api_key.strip() or not question.strip() or df is None))

    if run:
        if df is None:
            st.error("知識庫尚未就緒：請上傳有效 CSV 或放置預設檔案於同目錄。")
            return

        # 初始化向量資料庫（使用 cache）
        with st.spinner("正在初始化向量資料庫..."):
            collection = init_vector_db(df, api_key.strip())
        
        # 混合檢索
        with st.spinner("正在檢索相關段落（混合語義+關鍵字）..."):
            chunks = hybrid_retrieve_topk(collection, df, question, k=5)

        if not chunks:
            st.warning("找不到與問題相關的段落（以簡易字串比對）。你可以嘗試換個說法、加入更多關鍵字，或確認 CSV 的 `text_content` 內容。")
            return

        context = "\n\n---\n\n".join([c.text_content for c in chunks])
        prompt = build_prompt(context=context, question=question.strip())

        try:
            with st.spinner("正在呼叫 Gemini 1.5 Pro 生成回答 ..."):
                answer = generate_with_gemini(api_key.strip(), prompt)
        except Exception as e:
            st.error(str(e))
            return

        st.subheader("AI 回答")
        st.write(answer)

        with st.expander("📚 查看參考來源 / References"):
            for i, c in enumerate(chunks, start=1):
                # 構建 metadata 顯示
                page_info = f"第 {c.page_number} 頁" if c.page_number else ""
                year_info = f"{c.year} 年" if c.year else "N/A"
                
                st.markdown(
                    f"**#{i}** 📄 `{c.source_file}` {f'({page_info})' if page_info else ''}\n\n"
                    f"類型：`{c.doc_type}` ｜年份：`{year_info}` ｜相關度：`{c.score}%`"
                )
                st.write(c.text_content)
                st.divider()


if __name__ == "__main__":
    main()

